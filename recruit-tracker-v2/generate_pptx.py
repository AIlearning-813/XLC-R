#!/usr/bin/env python3
"""将阶段汇报内容生成 PPTX 文件"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== 颜色方案 =====
NAVY      = RGBColor(0x0A, 0x1F, 0x3F)
NAVY_LIGHT = RGBColor(0x14, 0x42, 0x72)
GOLD      = RGBColor(0xB8, 0x94, 0x4C)
GOLD_LIGHT = RGBColor(0xD4, 0xAF, 0x6C)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_50   = RGBColor(0xF8, 0xF9, 0xFB)
GRAY_100  = RGBColor(0xED, 0xF0, 0xF4)
GRAY_200  = RGBColor(0xD8, 0xDD, 0xE4)
GRAY_500  = RGBColor(0x7A, 0x84, 0x94)
GRAY_600  = RGBColor(0x5A, 0x64, 0x75)
GRAY_800  = RGBColor(0x2A, 0x33, 0x43)
GREEN     = RGBColor(0x48, 0xBB, 0x78)
RED       = RGBColor(0xD1, 0x52, 0x3F)
GREEN_DK  = RGBColor(0x2D, 0x7A, 0x3F)
BLUE_CARD = RGBColor(0x3B, 0x82, 0xC4)
GREEN_CARD= RGBColor(0x48, 0xBB, 0x78)

FONT_TITLE = '微软雅黑'
FONT_BODY  = '微软雅黑'
FONT_NUM   = 'Arial'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ===== 工具函数 =====
def add_textbox(slide, left, top, width, height, text="", font_size=12,
                font_name=FONT_BODY, bold=False, color=GRAY_800, alignment=PP_ALIGN.LEFT,
                line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf

def add_rich_textbox(slide, left, top, width, height):
    """返回 text_frame 供调用方添加多个段落"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_para(tf, text, font_size=12, font_name=FONT_BODY, bold=False, color=GRAY_800,
             alignment=PP_ALIGN.LEFT, space_after=0, first=False):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    return p

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    return shape

def add_round_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape

def add_line(slide, left, top, width, color=GOLD, thickness=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1: 封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# 深蓝背景（近似渐变用两个矩形叠加）
add_rect(slide1, 0, 0, W, H, fill_color=NAVY)

# 顶部 tag
add_textbox(slide1, Inches(5.5), Inches(1.0), Inches(2.3), Inches(0.4),
            "内部汇报 · 2026.07", font_size=9, font_name=FONT_NUM,
            bold=True, color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)

# 主标题
add_textbox(slide1, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
            "招聘管理系统", font_size=48, font_name=FONT_TITLE,
            bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# 副标题
add_textbox(slide1, Inches(2), Inches(3.0), Inches(9.3), Inches(0.6),
            "零开发基础 · 一个人 · 一个月 · 从想法到上线", font_size=18,
            color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

# highlight tag
tag_bg = add_round_rect(slide1, Inches(3.8), Inches(3.8), Inches(5.7), Inches(0.55),
                        fill_color=RGBColor(0x1A, 0x3F, 0x6E))
add_textbox(slide1, Inches(3.8), Inches(3.82), Inches(5.7), Inches(0.5),
            "需求方即产品经理 — 业务专家直接\"写\"出系统", font_size=13,
            color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)

# 分隔线
add_line(slide1, Inches(5.9), Inches(4.7), Inches(1.5), color=GOLD_LIGHT, thickness=Pt(1))

# 底部数字
stats_data = [("142", "候选人已入库"), ("9", "招聘专员在用"), ("✔", "已上线运行中")]
for i, (num, label) in enumerate(stats_data):
    x = Inches(3.3 + i * 2.5)
    add_textbox(slide1, x, Inches(5.1), Inches(2.0), Inches(0.8),
                num, font_size=36, font_name=FONT_NUM, bold=True,
                color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, x, Inches(5.8), Inches(2.0), Inches(0.4),
                label, font_size=10, color=RGBColor(0x99, 0x99, 0x99),
                alignment=PP_ALIGN.CENTER)


# ============================================================
# 内容页通用函数
# ============================================================
def make_content_slide(section_num, title):
    """创建白色背景内容页，返回 slide 对象"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill_color=WHITE)

    # 顶部色条
    add_rect(slide, 0, 0, W, Pt(4), fill_color=GOLD)

    # 章节号
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.3),
                section_num, font_size=8, font_name=FONT_NUM,
                bold=True, color=GOLD)

    # 标题 + 金色下划线
    add_textbox(slide, Inches(0.8), Inches(0.85), Inches(6), Inches(0.6),
                title, font_size=22, font_name=FONT_TITLE,
                bold=True, color=NAVY)
    add_line(slide, Inches(0.8), Inches(1.45), Inches(1.2), thickness=Pt(3))

    return slide


# ============================================================
# SLIDE 2: 成本分析
# ============================================================
slide2 = make_content_slide("01 / 成本分析", "做这样一套系统，市面上要花多少钱")

# 三列卡片
card_w = Inches(3.5)
card_h = Inches(2.6)
card_gap = Inches(0.4)
card_y = Inches(2.0)
start_x = Inches(0.8)

cost_data = [
    ("外包定制开发", "¥15~30万", "一次性交付，后续改需求\n另收费 ¥1,500~3,000/天", "ext"),
    ("SaaS 订阅", "¥6~15万", "每年都要付，功能固定\n无法按业务定制", "ext"),
    ("我们的方式", "≈ ¥0", "AI 辅助自主开发\n零外包零采购", "us"),
]

for i, (lbl, prc, sub, style) in enumerate(cost_data):
    x = start_x + i * (card_w + card_gap)
    bg_color = RGBColor(0xFF, 0xF8, 0xF0) if style == "ext" else RGBColor(0xF0, 0xF8, 0xF0)
    prc_color = RED if style == "ext" else GREEN_DK

    add_round_rect(slide2, x, card_y, card_w, card_h, fill_color=bg_color)

    add_textbox(slide2, x, card_y + Inches(0.4), card_w, Inches(0.3),
                lbl, font_size=10, color=GRAY_500, alignment=PP_ALIGN.CENTER)
    add_textbox(slide2, x, card_y + Inches(0.9), card_w, Inches(0.8),
                prc, font_size=32, font_name=FONT_NUM, bold=True,
                color=prc_color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide2, x + Inches(0.3), card_y + Inches(1.7), Inches(2.9), Inches(0.7),
                sub.replace('\n', '   '), font_size=8,
                color=GRAY_500, alignment=PP_ALIGN.CENTER)

# 底部备注
note_y = card_y + card_h + Inches(0.4)
note_bg = add_round_rect(slide2, start_x, note_y, Inches(11.7), Inches(1.3), fill_color=GRAY_50)
tf = add_rich_textbox(slide2, start_x + Inches(0.5), note_y + Inches(0.15), Inches(10.7), Inches(1.0))

add_para(tf, "省下至少 ¥15~30 万首次开发费，且每年省 ¥6~15 万订阅费", font_size=14,
         bold=True, color=GREEN_DK, alignment=PP_ALIGN.CENTER, first=True)
add_para(tf, "后期运维仅需云服务：约 ¥200~500/月（托管 + AI调用 + 存储）", font_size=11,
         color=GRAY_600, alignment=PP_ALIGN.CENTER)
add_para(tf, "对比外包运维报价 ¥2,000~5,000/月，每年再省 ¥2~5 万", font_size=11,
         color=GRAY_600, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: AI 驱动开发
# ============================================================
slide3 = make_content_slide("02 / 怎么做到的", "零编程基础，如何一个人把系统做出来")

steps = [
    ("１", "项目规划书，迭代打磨",
     "先画清楚\"系统要做什么\"——数据怎么存、流程怎么走、权限怎么分。不是写代码，是画蓝图。反复推翻重来，直到图纸上没有逻辑漏洞。"),
    ("２", "AI 编程助手（Claude Code）执行",
     "把规划书交给 AI，它来写代码。业务专家负责审——对不对、好不好用、逻辑通不通——AI 负责改。"),
    ("３", "对话即开发，反馈即迭代",
     "\"这里数据不对\"→ AI 查数据库定位 → 修复 → 验证。整个过程是人机对话，不需要写一行代码。"),
    ("４", "上线，持续优化",
     "系统部署到云端，专员开始用。发现问题→反馈→修复，迭代周期以小时计，不是周。"),
]

step_y = Inches(2.1)
for i, (n, title, desc) in enumerate(steps):
    y = step_y + i * Inches(1.2)
    # 左侧色条
    add_rect(slide3, Inches(0.8), y, Pt(3), Inches(1.0), fill_color=GOLD)
    # 背景
    add_rect(slide3, Inches(0.9), y, Inches(7.8), Inches(1.0), fill_color=GRAY_50)
    # 编号
    add_textbox(slide3, Inches(1.1), y + Inches(0.1), Inches(0.5), Inches(0.7),
                n, font_size=26, font_name=FONT_NUM, bold=True, color=GOLD)
    # 标题
    add_textbox(slide3, Inches(1.7), y + Inches(0.08), Inches(6.5), Inches(0.35),
                title, font_size=12, font_name=FONT_TITLE, bold=True, color=NAVY)
    # 描述
    add_textbox(slide3, Inches(1.7), y + Inches(0.45), Inches(6.8), Inches(0.5),
                desc, font_size=9, color=GRAY_500)

# 右侧深色大卡
r_x = Inches(9.2)
r_w = Inches(3.5)
r_h = Inches(4.85)
add_rect(slide3, r_x, step_y, r_w, r_h, fill_color=NAVY)

add_textbox(slide3, r_x, step_y + Inches(0.6), r_w, Inches(1.2),
            "0", font_size=72, font_name=FONT_NUM, bold=True,
            color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)
add_textbox(slide3, r_x, step_y + Inches(1.6), r_w, Inches(0.4),
            "行代码亲手写", font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC),
            alignment=PP_ALIGN.CENTER)

add_line(slide3, r_x + Inches(1.4), step_y + Inches(2.3), Inches(0.7),
         color=GOLD_LIGHT, thickness=Pt(1))

insight_text = (
    "真正的含金量不在代码里\n\n"
    "而在于：\n"
    "· 把业务经验变成系统设计\n"
    "· 知道什么是对的，什么是错的\n"
    "· 能持续对话、持续调整\n\n"
    "这不是买来的标准化产品\n"
    "这是业务 know-how 的数字化沉淀"
)
add_textbox(slide3, r_x + Inches(0.3), step_y + Inches(2.6), Inches(2.9), Inches(2.0),
            insight_text, font_size=9, color=RGBColor(0x99, 0x99, 0x99),
            alignment=PP_ALIGN.CENTER, line_spacing=1.5)


# ============================================================
# SLIDE 4: 功能一览
# ============================================================
slide4 = make_content_slide("03 / 系统功能", "目前能做什么")

features = [
    ("📬 邮箱自动收简历", "绑定招聘邮箱后定时扫描，简历附件自动入库，支持BOSS直聘、智联等平台", BLUE_CARD),
    ("🤖 AI 自动解析", "上传即自动提取姓名、学历、经历等关键信息，支持 PDF/Word/图片，准确率 >95%", GOLD),
    ("📋 流程看板", "筛选→初试→复试→终试→入职，拖拽流转，每个候选人进度一目了然", GREEN_CARD),
    ("📊 数据报表", "部门招聘进度、渠道转化效果、周期趋势，图表自动生成，一键导出", RGBColor(0xE0, 0x7B, 0x3C)),
    ("💬 AI 招聘助手", "基于公司知识库和历史数据，提供候选人匹配建议和招聘策略咨询", RGBColor(0x7B, 0x5E, 0xA7)),
    ("⚙️ 权限与安全", "分级权限、操作全留痕、每日自动备份、回收站防误删", RGBColor(0x3B, 0x8B, 0x9C)),
]

feat_w = Inches(3.6)
feat_h = Inches(2.0)
feat_gap_x = Inches(0.25)
feat_gap_y = Inches(0.3)
feat_start_x = Inches(0.8)
feat_start_y = Inches(2.0)

for i, (title, desc, border_color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = feat_start_x + col * (feat_w + feat_gap_x)
    y = feat_start_y + row * (feat_h + feat_gap_y)

    # 卡片背景
    add_rect(slide4, x, y, feat_w, feat_h, fill_color=GRAY_50)
    # 顶部色条
    add_rect(slide4, x, y, feat_w, Pt(4), fill_color=border_color)

    add_textbox(slide4, x + Inches(0.25), y + Inches(0.3), Inches(3.1), Inches(0.35),
                title, font_size=13, font_name=FONT_TITLE, bold=True, color=NAVY)
    add_textbox(slide4, x + Inches(0.25), y + Inches(0.8), Inches(3.1), Inches(1.0),
                desc, font_size=10, color=GRAY_600, line_spacing=1.5)


# ============================================================
# SLIDE 5: 开发历程
# ============================================================
slide5 = make_content_slide("04 / 开发历程", "一个月内做了什么")

# 四个数字卡片
num_data = [("151", "版本迭代", False), ("253", "需求与问题闭环", True),
            ("34", "安全隐患修复", False), ("0", "数据丢失", True)]
num_w = Inches(2.7)
num_gap = Inches(0.3)
num_start_x = Inches(0.8)
num_y = Inches(2.0)

for i, (num, label, is_gold) in enumerate(num_data):
    x = num_start_x + i * (num_w + num_gap)
    add_round_rect(slide5, x, num_y, num_w, Inches(1.6), fill_color=GRAY_50)
    c = GOLD if is_gold else NAVY
    add_textbox(slide5, x, num_y + Inches(0.15), num_w, Inches(0.8),
                num, font_size=36, font_name=FONT_NUM, bold=True,
                color=c, alignment=PP_ALIGN.CENTER)
    add_textbox(slide5, x, num_y + Inches(1.05), num_w, Inches(0.4),
                label, font_size=10, color=GRAY_500, alignment=PP_ALIGN.CENTER)

# 时间线
tl_y = num_y + Inches(2.0)
tl_left = Inches(1.3)
tl_right = Inches(12.0)
tl_w = tl_right - tl_left

# 横线
add_rect(slide5, tl_left, tl_y + Inches(0.15), tl_w, Pt(2), fill_color=GRAY_200)

tl_items = [
    ("6月上旬", "方案设计", "梳理业务流程\n确定系统框架"),
    ("6月中旬", "核心功能开发", "简历解析 · 邮箱归集\n流程看板 · 数据报表"),
    ("6月下旬", "安全 + 完善", "权限审计 · 隐私合规\nAI助手 · 回收站"),
    ("7月 →", "上线运营", "日常使用保障\n按需迭代优化"),
]

for i, (time_str, title, desc) in enumerate(tl_items):
    x = tl_left + Inches(0.5) + i * Inches(2.7)
    # 圆点
    dot = slide5.shapes.add_shape(MSO_SHAPE.OVAL, x, tl_y + Inches(0.05), Inches(0.22), Inches(0.22))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GOLD
    dot.line.fill.background()

    add_textbox(slide5, x - Inches(0.6), tl_y + Inches(0.4), Inches(1.8), Inches(0.25),
                time_str, font_size=9, font_name=FONT_NUM, bold=True,
                color=GOLD, alignment=PP_ALIGN.CENTER)
    add_textbox(slide5, x - Inches(0.6), tl_y + Inches(0.65), Inches(1.8), Inches(0.3),
                title, font_size=11, font_name=FONT_TITLE, bold=True,
                color=NAVY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide5, x - Inches(0.6), tl_y + Inches(0.95), Inches(1.8), Inches(0.5),
                desc, font_size=8, color=GRAY_500, alignment=PP_ALIGN.CENTER)

# PM tag
tag_x = Inches(3.8)
tag_y = tl_y + Inches(1.8)
tag_w = Inches(5.7)
tag_bg = add_round_rect(slide5, tag_x, tag_y, tag_w, Inches(0.45), fill_color=NAVY)
add_textbox(slide5, tag_x, tag_y + Inches(0.05), tag_w, Inches(0.35),
            "🎯 需求方 = 产品经理：不需要翻译需求，不需要等排期，想到就做",
            font_size=11, color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6: 现状与展望
# ============================================================
slide6 = make_content_slide("05 / 现状与展望", "已完成 & 下一步")

done_items = [
    "简历自动收录 + AI 解析",
    "候选人流程看板（拖拽）",
    "数据分析报表",
    "招聘需求申请与审批",
    "AI 招聘助手",
    "权限管理 · 操作留痕",
    "回收站 · 每日备份",
    "手机端可用",
]
next_items = [
    ("异常监控告警", "简历入库失败自动通知，第一时间处理", True),
    ("手机端体验优化", "移动审批、手机看简历更流畅", True),
    ("候选人智能推荐", "岗位与候选人自动匹配评分", False),
    ("用人部门协作", "部门负责人直接查看进度、反馈面试意见", False),
]

col_w = Inches(5.6)
col_gap = Inches(0.3)
col_y = Inches(2.0)
col_h = Inches(4.5)

for col_idx, (title, dot_color, items, is_list) in enumerate([
    ("已上线", GREEN, done_items, True),
    ("下一阶段", GOLD, next_items, False),
]):
    x = Inches(0.8) + col_idx * (col_w + col_gap)

    add_rect(slide6, x, col_y, col_w, col_h, fill_color=GRAY_50)

    # 标题行
    dot = slide6.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), col_y + Inches(0.4),
                                   Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()

    add_textbox(slide6, x + Inches(0.6), col_y + Inches(0.33), Inches(3), Inches(0.3),
                title, font_size=13, font_name=FONT_TITLE, bold=True, color=NAVY)

    if is_list:
        for j, item in enumerate(done_items):
            iy = col_y + Inches(0.9) + j * Inches(0.38)
            add_textbox(slide6, x + Inches(0.6), iy, Inches(4.5), Inches(0.35),
                        "✓  " + item, font_size=10, color=GRAY_600)
    else:
        for j, (strong, span, active) in enumerate(next_items):
            iy = col_y + Inches(0.9) + j * Inches(0.85)
            # 左边框
            border_c = BLUE_CARD if active else GRAY_200
            add_rect(slide6, x + Inches(0.4), iy, Pt(2), Inches(0.65), fill_color=border_c)
            txt_c = GRAY_800 if active else GRAY_500
            add_textbox(slide6, x + Inches(0.6), iy, Inches(4.5), Inches(0.3),
                        strong, font_size=11, font_name=FONT_TITLE, bold=True, color=txt_c)
            span_c = GRAY_500 if active else RGBColor(0xAA, 0xAA, 0xAA)
            add_textbox(slide6, x + Inches(0.6), iy + Inches(0.3), Inches(4.5), Inches(0.3),
                        span, font_size=9, color=span_c)


# ===== 保存 =====
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '阶段汇报.pptx')
prs.save(out_path)
print(f"PPTX 已生成: {out_path}")
